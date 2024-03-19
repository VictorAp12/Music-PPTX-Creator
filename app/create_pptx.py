"""
Este módulo tem como objetivo criar apresentações de slides em PowerPoint para musicas.
Autor: Victor
Data: setembro 2023 ~ dezembro 2023
"""

import os
import re
from dataclasses import dataclass
from typing import List, Literal, Tuple, Union
from io import BytesIO

import requests

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from PIL import Image

from PySide6.QtWidgets import QWidget

# pylint: disable=E0401,E1101, C0301, C0103, E1136, W0212


def _format_filename(name: str) -> str:
    """
    Convert a string to a valid filename.

    :param name (str): The file_name.

    :returns: The formatted file_name.
    """

    return re.sub(r"[\\|/^:*\"<>?\t]", "", name).strip()


@dataclass
class SlidesConfig:
    """
    This class represents the configuration for the slides.

    :param background_color_RGB: The color of the background in RGB (its a tuple).

    :param title_font_size_PT: The size of the title font. In points (its an integer).

    :param title_font_color_RGB: The color of the title font in RGB (its a tuple).

    :param title_font_isbold: If the title font is bold (its a boolean).

    :param title_font_name: The name of the title font.

    :param subtitle_font_size_PT: The size of the subtitle font. In points (its an integer).

    :param subtitle_font_color_RGB: The color of the subtitle font in RGB (its a tuple).

    :param subtitle_font_isbold: If the subtitle font is bold (its a boolean).

    :param subtitle_font_name: The name of the subtitle font.

    :param text_font_size_PT: The size of the text font. In points (its an integer).

    :param text_font_color_RGB: The color of the text font in RGB (its a tuple).

    :param text_font_isbold: If the text font is bold (its a boolean).

    :param text_font_alignment: The alignment of the text font. CENTER, LEFT, RIGHT, JUSTIFY and JUSTIFY_LOW.

    :param text_font_line_spacing_PT: The line spacing of the text font. In points (its an integer).

    :param text_font_space_before_PT: The space before the text font. In points (its an integer).

    :param text_font_space_after_PT: The space after the text font. In points (its an integer).

    :param text_font_name: The name of the text font.
    """

    background_color_RGB: Tuple = (255, 255, 255)
    title_font_size_PT: int = 56
    title_font_color_RGB: Tuple = (0, 0, 0)
    title_font_isbold: bool = False
    title_font_name: str = "Tw Cen MT"
    subtitle_font_size_PT: int = 44
    subtitle_font_color_RGB: Tuple = (0, 0, 0)
    subtitle_font_isbold: bool = False
    subtitle_font_name: str = "Tw Cen MT"
    text_font_size_PT: int = 38
    text_font_color_RGB: Tuple = (0, 0, 0)
    text_font_isbold: bool = False
    text_font_alignment: PP_PARAGRAPH_ALIGNMENT | int = PP_PARAGRAPH_ALIGNMENT.CENTER  # type: ignore
    text_font_line_spacing_PT: int = 38
    text_font_space_before_PT: int = 0
    text_font_space_after_PT: int = 0
    text_font_name: str = "Arial"


def create_slides(
    widget: QWidget,
    music_title: str,
    music_singer: str,
    music_lyric: List[str],
    slides_config: SlidesConfig = SlidesConfig(),
    FILES_FOLDER: str = "./files",
    method: str = "page_one",
    genius_image_link: str = "",
    image: bytes | None = None,
    background_image: bytes | BytesIO | None = None,
    language: Literal["pt", "en"] = "pt",
) -> Tuple[str, str, str]:
    """
    Function that creates presentation slides based on a song lyric.
    :param music_title (str): The title of the song.
    :param music_singer (str): The singer of the song.
    :param music_lyric (list[str]): The lyric of the song, formatted as a list of strings. Each item in the list must be a strophe.
    :param slides_config (SlidesConfig): The configuration for the slides.
    :param FILES_FOLDER (str, optional): The folder path. Defaults to "./files".
    :param method (Literal["page_one", "page_many", "page_insert_manually"], optional): The method to insert the lyrics. Defaults to "page_one".
    :param genius_image (str, optional): The api genius music image link. Defaults to ""
    :param image (bytes, optional): The image in bytes format. Defaults to None.
    :param background_image (bytes, optional): The background image in bytes format. Defaults to None.

    :return music (str): Music name, singer and the lyrics.
    :return file_name (str): The name of the file to be created.
    :return path (str): The path of the file to be created.
    """

    new_presentation = Presentation()

    # loop through each line (which is a strophe) in the lyric, creating a slide for each
    for i, paragraph_text in enumerate(music_lyric):
        new_slide = new_presentation.slides.add_slide(
            new_presentation.slide_layouts[0]
            if i == 0
            else new_presentation.slide_layouts[6]
        )

        # set the slide width and height to 16:9
        new_presentation.slide_width = Inches(16)
        new_presentation.slide_height = Inches(9)

        # background color
        background_color = new_slide.background
        fill = background_color.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*slides_config.background_color_RGB)

        if background_image:

            left = top = 0

            pic = new_slide.shapes.add_picture(
                background_image,
                left,
                top,
                width=new_presentation.slide_width,
                height=new_presentation.slide_height,
            )

            # set the image order to make it the content (over the rectangular box)
            new_slide.shapes._spTree.remove(pic._element)
            new_slide.shapes._spTree.insert(2, pic._element)

        if i == 0:
            title_text = music_title.upper()
            subtitle_text = music_singer.upper()

            # set the title and subtitle text for the new slide
            title_shape = new_slide.shapes.title
            title_shape.width = int(new_presentation.slide_width / 2)
            title_shape.height = int(new_presentation.slide_height / 5)
            title_shape.left = (
                int((new_presentation.slide_width - title_shape.width) / 2)
                if method == "page_insert_manually"
                else int(2)
            )
            title_shape.top = int(
                (new_presentation.slide_height - title_shape.height) / 4
            )
            title_shape.text = title_text

            subtitle_shape = new_slide.placeholders[1]
            subtitle_shape.width = int(new_presentation.slide_width / 2)
            subtitle_shape.height = int(new_presentation.slide_height / 5)
            subtitle_shape.left = (
                int((new_presentation.slide_width - title_shape.width) / 2)
                if method == "page_insert_manually"
                else int(2)
            )
            subtitle_shape.top = int((title_shape.top + title_shape.height))
            subtitle_shape.text = subtitle_text

            # configure the title format
            title_text_frame = title_shape.text_frame
            title_text_frame.paragraphs[0].font.color.rgb = RGBColor(
                *slides_config.title_font_color_RGB
            )
            title_text_frame.paragraphs[0].font.size = Pt(
                slides_config.title_font_size_PT
            )
            title_text_frame.paragraphs[0].font.name = slides_config.title_font_name
            title_text_frame.paragraphs[0].font.bold = slides_config.title_font_isbold

            # configure the subtitle format
            subtitle_text_frame = subtitle_shape.text_frame
            subtitle_text_frame.paragraphs[0].font.color.rgb = RGBColor(
                *slides_config.subtitle_font_color_RGB
            )
            subtitle_text_frame.paragraphs[0].font.size = Pt(
                slides_config.subtitle_font_size_PT
            )
            subtitle_text_frame.paragraphs[0].font.name = (
                slides_config.subtitle_font_name
            )
            subtitle_text_frame.paragraphs[0].font.bold = (
                slides_config.subtitle_font_isbold
            )

            if method in ("page_one", "page_many"):
                img_path = os.path.abspath(f"app/images/{music_title}.png")

                os.makedirs(os.path.abspath("app/images"), exist_ok=True)

                # verify if its an url
                if isinstance(genius_image_link, str) and genius_image_link.startswith(
                    "http"
                ):
                    response = requests.get(genius_image_link, timeout=5)
                    response.raise_for_status()
                    with open(img_path, "wb") as f:
                        f.write(response.content)

                    new_slide.shapes.add_picture(
                        img_path, Inches(9), Inches(2), Inches(6), Inches(4)
                    )

            if image:
                new_slide.shapes.add_picture(
                    Image.open(image), Inches(9), Inches(2), Inches(6), Inches(4)
                )

        # if it's not the first slide
        else:
            # set the position of the new slide
            x, y, width, height = (
                Inches(0),
                Inches(0),
                new_presentation.slide_width,
                new_presentation.slide_height,
            )

            # creating and formatting textbox in the new slide
            novo_textbox = new_slide.shapes.add_textbox(
                left=x, top=y, width=width, height=height
            ).text_frame

            new_paragraph = novo_textbox.add_paragraph()
            new_paragraph.font.name = slides_config.text_font_name
            new_paragraph.alignment = slides_config.text_font_alignment
            new_paragraph.font.bold = slides_config.text_font_isbold
            new_paragraph.font.color.rgb = RGBColor(*slides_config.text_font_color_RGB)
            new_paragraph.font.size = Pt(slides_config.text_font_size_PT)
            new_paragraph.line_spacing = Pt(
                slides_config.text_font_line_spacing_PT * 2
            )  # default = 1.2
            new_paragraph.space_before = Pt(slides_config.text_font_space_before_PT)
            new_paragraph.space_after = Pt(slides_config.text_font_space_after_PT)

            new_paragraph.text = paragraph_text.upper()

    file_name = _format_filename(
        music_title + " (" + music_singer.lower().strip() + ").pptx"
    )

    # saving the presentation
    new_presentation_path = os.path.join(FILES_FOLDER, file_name)

    if os.path.exists(new_presentation_path):
        os.remove(new_presentation_path)

    new_presentation.save(new_presentation_path)

    # if method == "page_one" or method == "page_insert_manually":
    if language == "pt":
        widget.setText(f'Arquivo: "{file_name}" concluído com sucesso!')  # type: ignore
    else:
        widget.setText(f'File: "{file_name}" completed!')  # type: ignore

    return (
        "{0} {1}\n{2}".format(music_title, music_singer, "\n\n".join(music_lyric)),
        file_name,
        str(os.path.join(FILES_FOLDER)),
    )


def extract_presentation_infos(
    presentation_path: str,
) -> Tuple[str, str, List[str], Union[bytes, None], Union[bytes, None]]:
    """
    Exctract the content from a presentation file.


    :param presentation_path: The path to the presentation file.
    :return: The content of the presentation which contains the title, the subtitle, the text,
        the image and the background image.
    """
    presentation = Presentation(presentation_path)

    text = [""]

    for i, slide in enumerate(presentation.slides):

        if i == 0:

            title = slide.shapes.title.text_frame.paragraphs[0].text
            subtitle = slide.placeholders[1].text

            image = None
            j = 1

            while j <= 3 and image is None:
                try:
                    image = slide.shapes[j].image.blob

                except (IndexError, AttributeError):
                    j += 1

            try:
                for item in slide.shapes:
                    if item.image:
                        background_image = item.image.blob
                        break

            except AttributeError:
                background_image = None

        else:

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame.text

                    text.append(text_frame.replace("\x0b", "\n"))

    return title, subtitle, text, image, background_image
